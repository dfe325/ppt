#prs = Presentation('test.pptx')
#prs.save('test.pptx')

#Source: https://stackoverflow.com/questions/39418620/extracting-text-from-multiple-powerpoint-files-using-python

'''from pptx import Presentation
import glob

for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)'''

#https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h

from pptx import Presentation
import glob


search_str = 'Purdue'
ppt = Presentation('EARTHRISE_ONE_SHEET_2.24.pptx')


li = []

for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and "E7119" in shape.text:
            print(shape.text)
            #li.append(shape.text)

print(li)

#for item in li:
#    print(type(item))




#iterate through slides in
