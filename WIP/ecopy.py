from docx import Document
import os

from pptx import Presentation
import glob

ppt = Presentation('EARTHRISE.pptx')

shape_li = []

for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and "E7120" in shape.text:
            #print(shape.text)
            shape_li.append(shape.text)


#print(shape_li)

document = Document()

document.add_paragraph(shape_li)

'''
li = []

for number in range(0, 3):
    li.append('new' + str(number))
'''
document.save('new.docx')

#print(li)

'''#creating new sequential list of docs
for item in li:
    document.save(item + '.docx')
'''
'''document.add_heading('Document Title', 0)

p = document.add_paragraph('A plain paragraph having some ')
p.add_run('bold').bold = True
p.add_run(' and some ')
p.add_run('italic.').italic = True
'''
document.save('ecopy.docx')
